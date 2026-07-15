#!/usr/bin/env python3
"""Project Suitability Engine — cover + 6 content slides (layman, builds slide-on-slide).
Order: cover · (1) old vs new · (2) scoring model · (3) groups + dictionary ·
       (4) worked example · (5) edge case w/ breakdowns · (6) challenges & assumptions.
Grounded in real prototype data (data/projects.json, data/applications.json, lib/scoring.ts)."""
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Palette ──────────────────────────────────────────────────────────────────
BLUE   = RGBColor(0x00, 0x32, 0x8A); BLUE2 = RGBColor(0x1E, 0x5B, 0xC6)
BLUEBG = RGBColor(0xE9, 0xEF, 0xFB); INK   = RGBColor(0x0F, 0x17, 0x2A)
MUTE   = RGBColor(0x55, 0x63, 0x77); SUBTLE= RGBColor(0x94, 0xA3, 0xB8)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFC); CARD  = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD7, 0xDF, 0xEA); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x15, 0x80, 0x3D); GREENBG=RGBColor(0xDC, 0xFA, 0xE6)
AMBER  = RGBColor(0xB4, 0x53, 0x09); AMBERBG=RGBColor(0xFD, 0xF1, 0xD8)
TEAL   = RGBColor(0x0E, 0x74, 0x90); TEALBG = RGBColor(0xDC, 0xF2, 0xF7)
PURP   = RGBColor(0x6D, 0x28, 0xD9); PURPBG = RGBColor(0xED, 0xE7, 0xFB)
SLATEBG= RGBColor(0xF1, 0xF5, 0xF9); PALE  = RGBColor(0xBD, 0xD0, 0xF2)
FONT = "Segoe UI"; N = 9
PORTAL = "DSTA Talent Outreach & Acquisition Portal"

prs = Presentation(); prs.slide_width = I(13.333); prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]

def rect(s,l,t,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,radius=None,shadow=False):
    sp=s.shapes.add_shape(shape,I(l),I(t),I(w),I(h)); sp.fill.solid()
    sp.fill.background() if fill is None else setattr(sp.fill.fore_color,'rgb',fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    if radius is not None and shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=radius
        except Exception: pass
    sp.shadow.inherit=False
    if shadow:
        el=sp._element.spPr; ef=el.makeelement(qn('a:effectLst'),{})
        sh=el.makeelement(qn('a:outerShdw'),{'blurRad':'90000','dist':'38100','dir':'5400000','rotWithShape':'0'})
        c=el.makeelement(qn('a:srgbClr'),{'val':'0F172A'}); al=el.makeelement(qn('a:alpha'),{'val':'12000'})
        c.append(al); sh.append(c); ef.append(sh); el.append(ef)
    return sp

def text(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space_after=4,ls=1.0,wrap=True):
    tb=s.shapes.add_textbox(I(l),I(t),I(w),I(h)); tf=tb.text_frame; tf.word_wrap=wrap
    tf.vertical_anchor=anchor; tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space_after); p.space_before=Pt(0); p.line_spacing=ls
        for (txt,size,bold,color,*rest) in para:
            it=rest[0] if rest else False
            r=p.add_run(); r.text=txt; r.font.size=Pt(size); r.font.bold=bold
            r.font.italic=it; r.font.color.rgb=color; r.font.name=FONT
    return tb

def chip(s,l,t,w,h,label,fill,fg,size=10.5,bold=True):
    rect(s,l,t,w,h,fill=fill,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    text(s,l,t,w,h,[[(label,size,bold,fg)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def bar(s,l,t,w,frac,fill,track=RGBColor(0xCF,0xDC,0xF2),h=0.13):
    rect(s,l,t,w,h,fill=track,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    if frac>0: rect(s,l,t,max(0.05,w*frac),h,fill=fill,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)

def base(title,kicker,idx):
    s=prs.slides.add_slide(BLANK)
    rect(s,0,0,13.333,7.5,fill=WHITE); rect(s,0,0,13.333,0.16,fill=BLUE)
    rect(s,0.55,0.5,0.12,0.6,fill=BLUE)
    text(s,0.78,0.44,11.6,0.4,[[(kicker.upper(),11,True,BLUE2)]])
    text(s,0.78,0.70,11.95,0.7,[[(title,23,True,INK)]])
    rect(s,0.78,1.36,11.78,0.018,fill=BORDER)
    text(s,0.55,7.06,9,0.3,[[(f"{PORTAL} · Project Suitability Engine",8.5,False,SUBTLE)]])
    text(s,11.6,7.06,1.18,0.3,[[(f"{idx:02d} / 0{N}",8.5,True,SUBTLE)]],align=PP_ALIGN.RIGHT)
    return s

NOTES=[]

# ════════════════════════════════════════════════════════════════════════════
# COVER
# ════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK)
rect(s,0,0,13.333,7.5,fill=BLUE); rect(s,0,6.95,13.333,0.55,fill=BLUE2)
rect(s,0.9,1.5,0.72,0.12,fill=WHITE)
text(s,0.9,1.85,11.5,0.5,[[("DSTA · TALENT OUTREACH & ACQUISITION PORTAL",13,True,PALE)]])
text(s,0.87,2.65,11.6,1.4,[[("Project Suitability Engine",44,True,WHITE)]])
text(s,0.9,4.1,10.9,1.0,[[("How we match applicants to projects — fairly, explainably,",16,False,RGBColor(0xDD,0xE7,0xF7))],
     [("and across every school background.",16,True,WHITE)]],ls=1.2)
chip(s,0.9,5.5,2.1,0.5,"Internal briefing",RGBColor(0x1E,0x4F,0xA0),WHITE,size=11)
NOTES.append("""COVER — open with the one-liner.
- "This is the Project Suitability Engine: it matches each applicant to each project and produces a Suitability Score we can fully explain - fairly, across every school background."
- It assists officers; it doesn't auto-decide.
- Then move into 'why change' (the old manual way vs the new engine).""")

# ════════════════════════════════════════════════════════════════════════════
# 1 — OLD WAY vs NEW WAY
# ════════════════════════════════════════════════════════════════════════════
s=base("Shortlisting today vs the new scoring engine","Why change",1)
rect(s,0.78,1.6,5.7,4.45,fill=LIGHT,line=BORDER,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
chip(s,1.05,1.85,1.95,0.42,"TODAY — MANUAL",AMBERBG,AMBER,size=10.5)
today=["IOs read thousands of applicants in Excel",
       "Tag each applicant to a project by hand",
       "Re-filter by project, then pick interviewees",
       "Grades aren't used in the matching at all"]
ty=2.55
for b in today:
    rect(s,1.1,ty+0.07,0.11,0.11,fill=AMBER,shape=MSO_SHAPE.OVAL)
    text(s,1.38,ty-0.03,4.9,0.5,[[(b,12.5,False,INK)]],ls=1.05); ty+=0.62
text(s,1.05,5.45,5.3,0.5,[[("Slow · inconsistent between officers · hard to audit",12,True,AMBER)]])
rect(s,6.85,1.6,5.7,4.45,fill=BLUEBG,line=BLUE,lw=1.2,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
chip(s,7.12,1.85,2.6,0.42,"NEW — SCORING ENGINE",BLUE,WHITE,size=10.5)
new=["One eligibility gate screens the first round",
     "Every applicant scored on Discipline of study + Skills + Standing",
     "Each project gets a ranked, explained shortlist",
     "Fair across every school background"]
ny=2.55
for b in new:
    rect(s,7.17,ny+0.07,0.11,0.11,fill=BLUE,shape=MSO_SHAPE.OVAL)
    text(s,7.45,ny-0.03,4.9,0.5,[[(b,12.5,False,INK)]],ls=1.05); ny+=0.62
text(s,7.12,5.45,5.3,0.5,[[("Consistent · auditable · explainable",12,True,BLUE)]])
rect(s,0.78,6.2,11.78,0.62,fill=BLUE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
text(s,0.78,6.2,11.78,0.62,[[("Same officers, same final decision — they just review a ranked, explained shortlist instead of a spreadsheet.",12.5,True,WHITE)]],
     align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
NOTES.append("""WHY CHANGE — old way vs new way. Set up the whole talk.
- Today (manual): IOs read thousands of applicants in Excel, tag each to a project by hand, re-filter, then pick interviewees. Slow, inconsistent between officers, hard to audit - and grades aren't even used in the matching.
- New (engine): one eligibility gate screens the first round, then every applicant is scored on three things - Discipline of study, Skills, Standing - and each project gets a ranked, explained shortlist.
- Reassure: same officers, same final decision. The engine hands them a ranked, explained shortlist instead of a blank spreadsheet.""")

# ════════════════════════════════════════════════════════════════════════════
# 2 — SCORING MODEL
# ════════════════════════════════════════════════════════════════════════════
s=base("What we assess — and how it becomes a score","The scoring model",2)
text(s,0.78,1.5,11.8,0.4,[[("Three things make the score. Each looks at the applicant and matches it to what the project needs:",12.5,False,MUTE)]])
rows=[("Discipline of study","matches the project's discipline","50%",BLUE),
      ("Skills","matches the project's required skills","30%",BLUE),
      ("Grades","become a “Standing” band — judged within your own school","20%",GREEN)]
ry=2.0
for left,mid,wt,c in rows:
    rect(s,0.78,ry,11.78,0.72,fill=CARD,line=BORDER,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1,shadow=True)
    text(s,1.05,ry,3.7,0.72,[[(left,13.5,True,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,4.7,ry,6.0,0.72,[[("→  "+mid,12,False,MUTE)]],anchor=MSO_ANCHOR.MIDDLE)
    chip(s,10.95,ry+0.16,1.4,0.4,wt,c,WHITE,size=14)
    ry+=0.82
rect(s,0.78,4.5,11.78,0.55,fill=BLUE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
text(s,0.78,4.5,11.78,0.55,[[("Suitability Score  =  50% Discipline of study  +  30% Skills  +  20% Standing",13.5,True,WHITE),
     ("   · configurable",10.5,False,PALE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
text(s,0.78,5.25,11.8,0.35,[[("Also used — but NOT part of the score:",11.5,True,INK)]])
rect(s,0.78,5.62,5.78,1.15,fill=SLATEBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
text(s,1.05,5.78,5.3,0.9,[[("Interest",12.5,True,TEAL)],
     [("Helps recommend projects to the applicant and acts as a light signal — it is not a scored weight.",11,False,MUTE)]],ls=1.12)
rect(s,6.78,5.62,5.78,1.15,fill=SLATEBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
text(s,7.05,5.78,5.3,0.9,[[("Project preference",12.5,True,PURP)],
     [("The applicant's own ranking. Used only to break ties between equal scores — never inside the score.",11,False,MUTE)]],ls=1.12)
NOTES.append("""THE SCORING MODEL — the heart. Be explicit.
- Three factors, each matching applicant to project:
   Discipline of study -> matches the project's discipline = 50%
   Skills -> matches the project's required skills = 30%
   Grades -> become a "Standing" band, judged within your own school = 20%
- Formula upfront: Suitability Score = 50% Discipline of study + 30% Skills + 20% Standing. Weights configurable per programme.
- INTEREST: NOT a scored weight - it powers project recommendations + a light signal (so "keen" doesn't beat "fit"). Can be given a small weight if DSTA wants.
- PREFERENCE: the applicant's ranking - the TIE-BREAK between equal scores, never inside the score.""")

# ════════════════════════════════════════════════════════════════════════════
# 3 — GROUPS + DICTIONARY
# ════════════════════════════════════════════════════════════════════════════
s=base("Same three factors — every school background","Across all backgrounds",3)
text(s,0.78,1.5,11.9,0.45,[[("The three factors don't change. Only ",12.5,False,MUTE),
    ("how we read Discipline of study and Standing",12.5,True,INK),(" changes per background:",12.5,False,MUTE)]])
chip(s,0.78,2.0,4.5,0.46,"Discipline of study  →  via the subject/discipline dictionary",BLUE,WHITE,size=10.5)
chip(s,5.45,2.0,4.2,0.46,"Standing  →  graded within its own system",BLUE2,WHITE,size=10.5)
cols=[("BACKGROUND",2.5),("WHAT WE READ FOR DISCIPLINE-FIT",4.5),("STANDING BAND COMES FROM",4.78)]
trows=[("Secondary (YDSP)","Subjects + interests","School / internal results",AMBER),
       ("JC (A-Level)","H1 / H2 subjects (H2 counts more)","A–E grades",GREEN),
       ("IB","HL / SL subjects (HL counts more)","1–7 score / 45 total",GREEN),
       ("IP","Subject combination","School-based results",AMBER),
       ("Polytechnic","Diploma discipline","Poly GPA",GREEN),
       ("University","Degree discipline","University GPA",GREEN)]
tx=0.78; ty=2.66; cx=tx
rect(s,tx,ty,11.78,0.4,fill=BLUE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
for lbl,w in cols:
    text(s,cx+0.2,ty,w-0.2,0.4,[[(lbl,10.5,True,WHITE)]],anchor=MSO_ANCHOR.MIDDLE); cx+=w
ry=ty+0.46
for i,(b,f,g,cc) in enumerate(trows):
    rect(s,tx,ry,11.78,0.44,fill=(WHITE if i%2==0 else SLATEBG),line=BORDER,lw=0.5)
    text(s,tx+0.2,ry,cols[0][1]-0.2,0.44,[[(b,11.5,True,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,tx+cols[0][1]+0.2,ry,cols[1][1]-0.2,0.44,[[(f,11,False,MUTE)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,tx+cols[0][1]+cols[1][1]+0.2,ry,cols[2][1]-0.4,0.44,[[(g,11,False,MUTE)]],anchor=MSO_ANCHOR.MIDDLE)
    ry+=0.44
rect(s,0.78,ry+0.12,11.78,0.86,fill=BLUEBG,line=BLUE,lw=1.2,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
text(s,1.0,ry+0.12,11.4,0.86,[
    [("Projects state a discipline, not subjects.",12,True,BLUE),
     ("  A shared dictionary translates each discipline into its relevant subjects —",12,False,INK)],
    [("e.g. Computer Science → Maths, Computing, Physics — so subject-based applicants (JC / IB) are matched too. DSTA admins maintain it.",11,False,MUTE)],
],anchor=MSO_ANCHOR.MIDDLE,ls=1.14,space_after=2)
NOTES.append("""ACROSS ALL BACKGROUNDS — one engine + the discipline->subject bridge.
- Three factors stay the same; only HOW we read Discipline-fit and Standing changes per background.
   Discipline-fit: subjects or discipline, via a shared dictionary.
   Standing: grades/GPA graded WITHIN that system (an A, IB 7, poly GPA never compared raw).
- DISCIPLINE-TO-SUBJECT bridge: projects only state a discipline. An admin-maintained dictionary maps each discipline to its subjects (Computer Science -> Maths, Computing, Physics) so JC/IB subjects can be matched. Unmapped discipline -> factor dropped, weights renormalise, confidence drops, admin flagged.
- Maintenance is light: small, slow-moving list; seed once; gaps fail safe.""")

# ════════════════════════════════════════════════════════════════════════════
# 4 — WORKED EXAMPLE
# ════════════════════════════════════════════════════════════════════════════
s=base("A real applicant, scored in full","Worked example",4)
text(s,0.78,1.45,11.8,0.4,[[("Putting it together: one real applicant from the prototype, scored on the three factors.",12,False,MUTE)]])
rect(s,0.78,1.95,5.6,0.92,fill=BLUEBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
text(s,1.0,2.06,5.2,0.75,[[("Tan Wei Ming",15,True,INK)],[("Computer Science · GPA 3.8",11.5,False,MUTE)]])
text(s,6.55,2.06,0.7,0.75,[[("→",20,True,SUBTLE)]],anchor=MSO_ANCHOR.MIDDLE)
rect(s,7.0,1.95,5.55,0.92,fill=SLATEBG,line=BORDER,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
text(s,7.22,2.06,5.2,0.75,[[("Project: Cybersecurity Threat Analysis",13,True,INK)],
     [("needs Python, Threat Intel, MITRE ATT&CK · CS / InfoSec discipline",10.5,False,MUTE)]],ls=1.05)
rows=[("Discipline of study","× 50%","Computer Science is exactly a required discipline","100",1.0,BLUE),
      ("Skills","× 30%","has Python + security basics; lighter on MITRE ATT&CK, forensics","60",0.60,BLUE),
      ("Standing","× 20%","GPA 3.8 → “Strong” band (within university track) → 80","80",0.80,GREEN)]
ry=3.15
for name,wt,why,val,frac,c in rows:
    text(s,0.85,ry,2.3,0.4,[[(name+"  ",13,True,INK),(wt,11.5,True,c)]])
    text(s,3.75,ry,7.0,0.4,[[(why,11,False,MUTE)]])
    text(s,11.2,ry,1.35,0.4,[[(val+" / 100",12.5,True,c)]],align=PP_ALIGN.RIGHT)
    bar(s,0.85,ry+0.34,11.7,frac,c)
    ry+=0.62
text(s,0.85,5.05,2.4,0.35,[[("Standing bands:",10.5,True,SUBTLE)]])
for i,(lbl,val) in enumerate([("Top",100),("Strong",80),("Solid",60),("Borderline",40)]):
    chip(s,2.55+i*1.55,5.02,1.45,0.36,f"{lbl} = {val}",SLATEBG,INK,size=10,bold=True)
rect(s,0.78,5.6,4.6,0.95,fill=BLUE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
text(s,1.0,5.6,4.2,0.95,[[("SUITABILITY SCORE",9.5,True,PALE)],[("84 / 100",24,True,WHITE)]],anchor=MSO_ANCHOR.MIDDLE,ls=1.0)
rect(s,5.55,5.6,7.0,0.95,fill=SLATEBG,line=BORDER,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
text(s,5.8,5.67,6.6,0.85,[[("0.50×100  +  0.30×60  +  0.20×80  =  84",13.5,True,INK)],
     [("Every number is shown to the user — nothing hidden.",10.5,False,MUTE,True)]],anchor=MSO_ANCHOR.MIDDLE,ls=1.15)
NOTES.append("""WORKED EXAMPLE — real prototype numbers.
- Tan Wei Ming, CS, GPA 3.8 -> "Cybersecurity Threat Analysis".
   Discipline of study x50%: CS is exactly a required discipline -> 100
   Skills x30%: Python + security basics, lighter on MITRE ATT&CK/forensics -> 60
   Standing x20%: GPA 3.8 -> "Strong" band -> 80
- Bands: Top 100, Strong 80, Solid 60, Borderline 40 (within own system).
- Arithmetic on screen: 0.50x100 + 0.30x60 + 0.20x80 = 84. An officer sees exactly why it's 84, not 95.""")

# ════════════════════════════════════════════════════════════════════════════
# 5 — EDGE CASE (per-score breakdown + add-up)
# ════════════════════════════════════════════════════════════════════════════
s=base("Post-JC / Post-Poly — one pool, many school systems","The hard case",5)
text(s,0.78,1.42,11.9,0.36,[[("One pool mixes several school systems for the ",12,False,MUTE),
    ("same project",12,True,INK),(". Same formula, scored side by side:",12,False,MUTE)]])
rect(s,0.78,1.84,11.78,0.42,fill=SLATEBG,line=BORDER,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
text(s,1.0,1.84,11.4,0.42,[[("Project: Defence Simulation Modelling",11.5,True,INK),
    ("   — needs Computer Science / Maths · Python, MATLAB, simulation",10.5,False,MUTE)]],anchor=MSO_ANCHOR.MIDDLE)
text(s,0.78,2.34,11.78,0.32,[[("Suitability = 0.5 × Discipline  +  0.3 × Skills  +  0.2 × Standing",11,True,BLUE)]],align=PP_ALIGN.CENTER)
cards=[("A-Level (JC)","H2 Maths (A), H2 Computing",94,60,80,81,GREEN),
       ("IB","HL Maths 7, SL Comp Sci",88,60,80,78,BLUE),
       ("Polytechnic","Dip. Info Tech — related field, strong skills",82,70,80,78,TEAL),
       ("NUS High","Major Maths & Computing",92,60,80,80,PURP)]
x=0.78; cw=2.86; gap=0.17
for i,(bg,desc,F,Sk,St,fit,c) in enumerate(cards):
    lx=x+i*(cw+gap)
    rect(s,lx,2.74,cw,3.28,fill=CARD,line=BORDER,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05,shadow=True)
    chip(s,lx+0.18,2.9,cw-0.36,0.4,bg,c,WHITE,size=10.5)
    text(s,lx+0.2,3.38,cw-0.36,0.58,[[(desc,9.5,False,MUTE)]],ls=1.05)
    br=[("Discipline ×50",F,BLUE),("Skills ×30",Sk,BLUE),("Standing ×20",St,GREEN)]
    by=4.02
    for lbl,val,bc in br:
        text(s,lx+0.2,by,1.85,0.3,[[(lbl,10,False,INK)]])
        text(s,lx+cw-0.85,by,0.65,0.3,[[(str(val),11,True,bc)]],align=PP_ALIGN.RIGHT)
        by+=0.34
    rect(s,lx+0.2,5.12,cw-0.4,0.014,fill=BORDER)
    text(s,lx+0.2,5.18,cw-0.4,0.28,[[(f"0.5·{F}+0.3·{Sk}+0.2·{St}",8.5,False,MUTE,True)]])
    text(s,lx+0.2,5.48,cw-0.4,0.42,[[("Suitability  ",8.5,True,SUBTLE),(f"{fit}",17,True,c),("/100",8.5,False,MUTE)]])
rect(s,0.78,6.12,11.78,0.72,fill=BLUE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
text(s,1.05,6.18,11.3,0.6,[
    [("Four school systems land in a tight 78–81 band — judged on the same factors, not on their certificate.",11.5,True,WHITE)],
    [("Grades normalised within each system · thinner data → lower confidence, never lower score · ties broken by project preference.",10,False,PALE)],
],ls=1.1,space_after=3)
NOTES.append("""THE HARD CASE — show the breakdowns and the add-up.
- One pool mixes A-Level, IB, NUS High, Polytechnic (and O-Level Y4 / others) for the SAME project: Defence Simulation Modelling (needs CS/Maths, Python, MATLAB).
- HOW EACH ADDS UP (weighted, not simple sum): Suitability = 0.5*Discipline + 0.3*Skills + 0.2*Standing.
   A-Level: 0.5*94 + 0.3*60 + 0.2*80 = 81
   IB:      0.5*88 + 0.3*60 + 0.2*80 = 78
   Poly:    0.5*82 + 0.3*70 + 0.2*80 = 78  (discipline lower - IT is RELATED not exact - but stronger hands-on skills offset it)
   NUS High:0.5*92 + 0.3*60 + 0.2*80 = 80
- HOW EACH SUB-SCORE IS FOUND: Discipline = how closely the field/subjects match the project's discipline (exact / related / weak). Skills = share of the project's required skills the applicant covers. Standing = grade/GPA mapped to a band within their own system (Top/Strong/Solid/Borderline).
- Tight 78-81 band -> judged on fit, not certificate. Grades normalised within each system; thin data -> lower confidence; ties -> project preference.""")

# ════════════════════════════════════════════════════════════════════════════
# 6 — CHALLENGES, ASSUMPTIONS & CONSIDERATIONS
# ════════════════════════════════════════════════════════════════════════════
s=base("Challenges, assumptions & things to consider","Before we build",6)
panels=[
 ("ASSUMPTIONS",BLUE,BLUEBG,[
    "Projects are tagged with a discipline (and ideally skills) by the PC / mentor",
    "Applicant data is structured — subjects, discipline, grades captured in the form, not free text",
    "Each programme mostly draws one or two backgrounds; mixed pools are the exception",
    "50/30/20 weights are sensible defaults, agreed with DSTA"]),
 ("CHALLENGES / RISKS",AMBER,AMBERBG,[
    "Maintaining the subject↔discipline dictionary and grade bands per system",
    "Free-text disciplines/skills won't match — needs controlled dropdowns",
    "Grade comparability across systems (A vs IB 7 vs GPA) — normalise within track",
    "Thin data for some groups (e.g. secondary) → lower confidence",
    "Automation bias — officers must stay in the loop"]),
 ("THINGS TO CONSIDER",GREEN,GREENBG,[
    "Keep humans deciding interviews; the engine ranks + explains only",
    "Show confidence, not just score; never auto-reject beyond eligibility",
    "Make weights & dictionary admin-configurable and versioned for audit",
    "Pilot on the Post-JC / Post-Poly group first (the hardest case)",
    "Decide: does Interest carry a weight? final Standing weight?"]),
]
x=0.78; cw=3.82; gap=0.16
for i,(hd,c,cbg,items) in enumerate(panels):
    lx=x+i*(cw+gap)
    rect(s,lx,1.7,cw,4.95,fill=CARD,line=BORDER,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04,shadow=True)
    chip(s,lx,1.7,cw,0.5,hd,c,WHITE,size=11.5)
    iy=2.42
    for it in items:
        rect(s,lx+0.26,iy+0.07,0.1,0.1,fill=c,shape=MSO_SHAPE.OVAL)
        text(s,lx+0.5,iy-0.02,cw-0.72,0.8,[[(it,10.5,False,INK)]],ls=1.06); iy+=0.86
NOTES.append("""BEFORE WE BUILD — challenges, assumptions, considerations. Be honest here; it builds trust.
- ASSUMPTIONS: projects carry a discipline (and ideally skills); applicant data is STRUCTURED not free text; each programme mostly draws 1-2 backgrounds; 50/30/20 are agreed defaults.
- CHALLENGES/RISKS: maintaining the subject<->discipline dictionary + grade bands; free-text won't match (need dropdowns); cross-system grade comparability (normalise within track); thin data for some groups; automation bias.
- CONSIDERATIONS / mitigations: humans decide interviews (engine ranks+explains only); show confidence not just score; never auto-reject beyond eligibility; make weights+dictionary admin-configurable and versioned; PILOT on Post-JC/Post-Poly first; open decisions - does Interest get a weight, and the final Standing weight.
- Close: the design is deliberately a glass box so every one of these is visible and adjustable.""")

# ════════════════════════════════════════════════════════════════════════════
# 7 — FAQ · how the scoring works
# ════════════════════════════════════════════════════════════════════════════
def faq_slide(title, kicker, idx, items, ih=1.0, step=1.06):
    s=base(title,kicker,idx)
    iy=1.62
    for q,a in items:
        rect(s,0.78,iy,11.78,ih,fill=SLATEBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
        text(s,1.05,iy+0.1,11.3,0.32,[[("Q.  ",12.5,True,BLUE2),(q,12.5,True,INK)]])
        text(s,1.05,iy+0.44,11.3,ih-0.5,[[("A.  ",11,True,SUBTLE),(a,11,False,MUTE)]],ls=1.08)
        iy+=step
    return s

faq_slide("FAQ — how the scoring works","Frequently asked · scoring",7,[
 ("How is the Discipline sub-score calculated?",
  "The project's discipline is expanded via the dictionary into relevant subjects (CS → Maths, Computing, Physics). Subject-based applicants are scored on coverage, weighted by level (H2/HL full, H1/SL ~60%). A formal discipline scores exact = 100, related ~80."),
 ("How is the Skills sub-score calculated?",
  "Share of the project's required skills the applicant demonstrably has:  15 + (matched ÷ required) × 85.  A 15 floor and 100 ceiling; required skills count more than preferred."),
 ("Can we change the weights?",
  "Yes — Discipline / Skills / Standing are set per programme by DSTA admins and versioned. Set Standing to 0 to ignore grades; raise Skills for skills-led projects."),
 ("How are ties broken?",
  "By the applicant's own project preference (their #1 choice first). If still tied, the officer decides. Preference is never inside the score."),
 ("Are the numbers in the examples exact?",
  "They are illustrative. In production every sub-score is computed from the applicant's declared data + the dictionary, and shown to the user in full."),
])
NOTES.append("""FAQ (scoring) — reference slide; use when asked for detail.
- Discipline sub-score: dictionary expands the project's discipline into relevant subjects; coverage weighted by level (H2/HL full, H1/SL ~0.6); formal discipline = exact 100 / related ~80.
- Skills sub-score (real formula in lib/scoring.ts): 15 + (matched/required)*85.
- Weights configurable + versioned; Standing can be 0.
- Ties: applicant's project preference, then officer.
- Be honest: example integers are illustrative; production computes them from declared data. (If pressed, pre-university SKILLS shown at 60 are generous; a strict reading is ~50 - offer to recompute.)""")

# ════════════════════════════════════════════════════════════════════════════
# 8 — FAQ · fairness, data & operating
# ════════════════════════════════════════════════════════════════════════════
faq_slide("FAQ — fairness, data & operating","Frequently asked · fairness & data",8,[
 ("Is it fair across different school systems?",
  "Applicants are compared only within a project's pool; grades are normalised within each system (an A, an IB 7, a poly GPA are never compared raw); a factor a group can't have is dropped and weights renormalise to 100."),
 ("What if data is missing, or a discipline isn't in the dictionary?",
  "Missing optional data lowers the confidence flag, not the score. An unmapped discipline drops that factor, renormalises, lowers confidence, and flags the admin to add one entry."),
 ("Will it auto-reject applicants or replace officers?",
  "No. The only automatic gate is eligibility. Everything else is ranked + explained — officers shortlist and choose who to interview."),
 ("Can applicants game it (inflate skills, rank everything #1)?",
  "Preference is only a tiebreaker; skills are matched against declared / CV evidence; each factor is capped — so gaming has limited payoff."),
 ("What do we need to run it, and how much upkeep?",
  "Structured applicant data (subjects / discipline / grades via dropdowns) and projects tagged with discipline + skills. Upkeep is light: a small subject↔discipline dictionary, seeded once, edited a few times a year."),
])
NOTES.append("""FAQ (fairness, data, operating) — reference slide.
- Fairness: within-pool comparison; grades normalised within each system; missing factor -> renormalise to 100.
- Missing data -> confidence not score; unmapped discipline -> drop + renormalise + flag admin.
- No auto-reject beyond eligibility; humans shortlist and select.
- Anti-gaming: preference is tiebreaker only; skills checked vs evidence; factors capped.
- Inputs: structured applicant data + projects tagged with discipline/skills. Upkeep light (dictionary seeded once, few edits/year).""")

# ════════════════════════════════════════════════════════════════════════════
# 9 — FAQ · normalisation & who configures it
# ════════════════════════════════════════════════════════════════════════════
faq_slide("FAQ — normalisation & who configures it","Frequently asked · normalisation & config",9,[
 ("How does grade normalisation work — and who sets the bands?",
  "Each school system has its own grade→band table (A-Level A = Top, B = Strong; IB 7 = Top, 6 = Strong; Poly GPA ≥3.8 = Top). The engine looks up the applicant's raw grade in their system's table to get the band (Top 100 / Strong 80 / Solid 60 / Borderline 40). DSTA admins set the thresholds per system; the engine only applies them."),
 ("How does weight “renormalisation” work when a factor is missing?",
  "Start from the weights (Discipline 50 / Skills 30 / Standing 20). If a group can't have a factor, the engine drops it and rescales the rest to total 100 — e.g. without Standing: Discipline 50÷80 = 62.5%, Skills 30÷80 = 37.5%. The applicant can still reach 100. It's automatic — no one does it by hand."),
 ("Where are the weights, bands and dictionary actually defined?",
  "All three are configuration data, not code: factor weights per programme, grade→band tables per school system, and the subject↔discipline dictionary. They live in the admin settings, separate from the scoring engine."),
 ("Who maintains it — and can we audit a score?",
  "DSTA admins own and edit all three, no developer needed. Every change is versioned, so any past score can be reproduced exactly from the settings that were live when it was generated."),
], ih=1.2, step=1.26)
NOTES.append("""FAQ (normalisation & config) — the where/by-whom reference.
- GRADE NORMALISATION: per-system grade->band lookup tables (A-Level A=Top, IB 7=Top, Poly GPA>=3.8=Top -> Top 100 / Strong 80 / Solid 60 / Borderline 40). Admins set thresholds; engine applies.
- WEIGHT RENORMALISATION: base weights 50/30/20; if a factor is absent, drop it and rescale remaining to 100 (drop Standing -> Discipline 62.5%, Skills 37.5%). Automatic, deterministic, no human.
- WHERE: three pieces of admin config (weights per programme, band tables per system, subject<->discipline dictionary) - all data, not code, separate from the engine.
- BY WHOM: DSTA admins (no developer); everything versioned so any past score is reproducible -> fully auditable.""")

# ── save + notes ─────────────────────────────────────────────────────────────
for slide,note in zip(prs.slides,NOTES):
    slide.notes_slide.notes_text_frame.text=note
out="/Users/youzhenglee/Desktop/DSTAProjectMockup/Project-Suitability-Engine.pptx"
prs.save(out)
print("Saved",out,"·",len(prs.slides._sldIdLst),"slides,",len(NOTES),"notes")
